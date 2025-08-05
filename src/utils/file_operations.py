import os
import fitz  # PyMuPDF
import openpyxl
from pptx import Presentation
import csv
import docx
import chardet
from extract_msg import Message
import email # 追加
from email import policy # 追加

# --- Text Extraction Functions ---
def extract_text_preview(filepath):
    """Extracts text from various file types for preview."""
    print(f"extract: {filepath}")
    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext == ".pdf":
            return extract_pdf_text(filepath)
        elif ext in [".xlsx", ".xlsm"]:
            return extract_excel_text(filepath)
        elif ext in [".pptx", ".pptm"]:
            return extract_pptx_text(filepath)
        elif ext in [".docx", ".docm"]:
            return extract_docx_text(filepath)
        elif ext == ".csv":
            return extract_csv_text(filepath)
        elif ext == ".msg":
            return extract_msg_text(filepath)
        elif ext == ".eml": # 追加
            return extract_eml_text(filepath) # 追加
        else:
            # For other extensions, attempt to read as a text file.
            return extract_text_file(filepath)
    except Exception as e:
        # This is a fallback for binary files or read errors.
        return f"プレビュー中にエラーが発生しました: {e}"

def extract_pdf_text(filepath):
    with fitz.open(filepath) as doc:
        return "".join(page.get_text() for page in doc)

def render_pdf_as_pixmaps(filepath, dpi=96):
    try:
        with fitz.open(filepath) as doc:
            matrix = fitz.Matrix(dpi / 72, dpi / 72)
            return [page.get_pixmap(matrix=matrix) for page in doc]
    except Exception as e:
        print(f"PDFレンダリングエラー: {e}")
        return []

def extract_excel_text(filepath):
    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    text = []
    for sheet in wb.worksheets:
        text.append(f"[{sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            line = "\t".join(str(cell) if cell is not None else "" for cell in row)
            text.append(line)
    return "\n".join(text)

def extract_pptx_text(filepath):
    prs = Presentation(filepath)
    text = []
    for i, slide in enumerate(prs.slides):
        text.append(f"[スライド{i + 1}]")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_docx_text(filepath):
    doc = docx.Document(filepath)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_csv_text(filepath):
    encoding = detect_encoding(filepath) or 'utf-8'
    rows = []
    try:
        with open(filepath, newline="", encoding=encoding, errors='ignore') as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i >= 200: # Limit rows for performance
                    break
                rows.append(", ".join(row))
        return "\n".join(rows)
    except (UnicodeDecodeError, csv.Error):
        # Fallback to raw text read if CSV parsing fails
        return extract_text_file(filepath)

def extract_text_file(filepath):
    """Reads a plain text file with robust encoding detection."""
    encoding = detect_encoding(filepath)
    try:
        with open(filepath, "r", encoding=encoding or 'utf-8', errors='ignore') as f:
            return f.read()
    except Exception:
        # If chardet fails or file is binary, read with a safe fallback.
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

def detect_encoding(filepath, sample_size=4096):
    """Detects the encoding of a file by reading a sample."""
    try:
        with open(filepath, 'rb') as f:
            sample = f.read(sample_size)
            result = chardet.detect(sample)
            return result['encoding']
    except (IOError, IndexError):
        return None

def extract_msg_text(filepath):
    """Extracts text content from .msg files."""
    try:
        msg = Message(filepath)
        text_content = []
        if msg.subject:
            text_content.append(f"Subject: {msg.subject}")
        if msg.sender:
            text_content.append(f"From: {msg.sender}")
        if msg.to:
            text_content.append(f"To: {msg.to}")
        if msg.cc:
            text_content.append(f"CC: {msg.cc}")
        if msg.date:
            text_content.append(f"Date: {msg.date}")
        text_content.append("---")
        if msg.body:
            text_content.append(msg.body)
        return "".join(text_content)
    except Exception as e:
        return f"Error extracting MSG file: {e}"

def extract_eml_text(filepath): # 追加
    """Extracts text content from .eml files."""
    try:
        with open(filepath, 'rb') as fp:
            msg = email.message_from_binary_file(fp, policy=policy.default)

        text_content = []
        if msg['subject']:
            text_content.append(f"Subject: {msg['subject']}")
        if msg['from']:
            text_content.append(f"From: {msg['from']}")
        if msg['to']:
            text_content.append(f"To: {msg['to']}")
        if msg['cc']:
            text_content.append(f"CC: {msg['cc']}")
        if msg['date']:
            text_content.append(f"Date: {msg['date']}")
        text_content.append("---")

        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type()
                cdispo = part.get('Content-Disposition')

                # extract plain text body
                if ctype == 'text/plain' and 'attachment' not in (cdispo or ''):
                    payload = part.get_payload(decode=True)
                    charset = part.get_content_charset()
                    if charset:
                        text_content.append(payload.decode(charset, errors='ignore'))
                    else:
                        text_content.append(payload.decode('utf-8', errors='ignore')) # Fallback
                    break # Only get the first plain text part
        else:
            payload = msg.get_payload(decode=True)
            charset = msg.get_content_charset()
            if charset:
                text_content.append(payload.decode(charset, errors='ignore'))
            else:
                text_content.append(payload.decode('utf-8', errors='ignore')) # Fallback

        return "\n".join(text_content)
    except Exception as e:
        return f"Error extracting EML file: {e}"